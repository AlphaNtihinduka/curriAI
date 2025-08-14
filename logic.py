
# logic.py (feature-toggled, end-to-end)
import os
import json
import time
import uuid
import zipfile
from typing import List, Optional, Dict, Any

from fastapi import FastAPI, APIRouter, UploadFile, File, Form, HTTPException
from fastapi.responses import FileResponse
import pdfplumber
import docx
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openai import OpenAI
from dotenv import load_dotenv

# PPT/PDF helpers
from pptx import Presentation
from pptx.util import Inches as PPTInches
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN

# PDF export path A: docx2pdf (Windows/macOS)
DOCX2PDF_AVAILABLE = False
try:
    from docx2pdf import convert as docx2pdf_convert
    DOCX2PDF_AVAILABLE = True
except Exception:
    pass

# PDF export path B: reportlab fallback (simple text-only)
REPORTLAB_AVAILABLE = False
try:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
    REPORTLAB_AVAILABLE = True
except Exception:
    pass

# -----------------------------
# Setup
# -----------------------------
load_dotenv()
# IMPORTANT: set OPENAI_API_KEY in your environment, or replace "<REPLACE_ME>"
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY") or "<REPLACE_ME>"
client = OpenAI(api_key=OPENAI_API_KEY)

app = FastAPI(title="Course Assets Generator")
router = APIRouter()

# -----------------------------
# Small helpers
# -----------------------------
def to_bool(v: str) -> bool:
    """Parse various stringy truthy values from form-data."""
    return str(v).strip().lower() in {"1", "true", "yes", "y", "on"}

# -----------------------------
# Utils: text extraction
# -----------------------------
def extract_text_from_file(file_path: str) -> str:
    if file_path.lower().endswith(".pdf"):
        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                t = page.extract_text() or ""
                if t.strip():
                    text_parts.append(t)
        return "\n".join(text_parts)
    elif file_path.lower().endswith(".docx"):
        d = docx.Document(file_path)
        return "\n".join(p.text for p in d.paragraphs if p.text)
    else:
        raise ValueError("Unsupported file type. Only PDF and DOCX are allowed.")

# -----------------------------
# LLM helpers
# -----------------------------
def llm_json(prompt: str, model: str = "gpt-4o-mini", temperature: float = 0.3) -> Dict[str, Any]:
    """Ask the model to return STRICT JSON. Validate with json.loads."""
    messages = [
        {"role": "system", "content": "Return ONLY valid JSON. No backticks, no prose."},
        {"role": "user", "content": prompt},
    ]
    resp = client.chat.completions.create(
        model=model,
        messages=messages,
        temperature=temperature,
        response_format={"type": "json_object"},
        timeout=60.0
    )
    content = resp.choices[0].message.content
    try:
        return json.loads(content)
    except Exception:
        content = content.strip().removeprefix("```json").removesuffix("```").strip()
        return json.loads(content)

# -----------------------------
# Planning (LLM): weeks, diagram hints, question blueprints
# -----------------------------
def plan_course_structure(descriptor_text: str, instructions: str, weeks: int) -> dict:
    prompt = (
        f"You are a senior curriculum designer. Create a {weeks}-week course plan with this structure:\n\n"
        "{\n"
        '  "course_meta": {\n'
        '    "title": "",\n'
        '    "code": "",\n'
        '    "credits": 3\n'
        "  },\n"
        '  "weeks": [\n'
        "    {\n"
        '      "week": 1,\n'
        '      "title": "",\n'
        '      "topics": [],\n'
        '      "objectives": [],\n'
        '      "activities": [],\n'
        '      "assessments": [],\n'
        '      "source_queries": [],\n'
        '      "diagram_suggestions": []\n'
        "    }\n"
        "  ],\n"
        '  "question_blueprints": {\n'
        '    "mcq": 8,\n'
        '    "short_answer": 4,\n'
        '    "essay": 2\n'
        "  }\n"
        "}\n\n"
        "Rules:\n"
        "- Keep topics practical and specific.\n"
        "- Include diagram suggestions only when they fit the content (flowchart/pie/bar).\n"
        "- Do NOT invent citations; only provide search queries.\n"
        "- Respect the lecturer instructions.\n\n"
        "Lecturer instructions:\n"
        "<<INSTRUCTIONS>>\n"
        f"{instructions.strip()}\n"
        "<</INSTRUCTIONS>>\n\n"
        "Course Descriptor:\n"
        "<<DESCRIPTOR>>\n"
        f"{descriptor_text[:30000]}\n"
        "<</DESCRIPTOR>>\n"
    )
    return llm_json(prompt)


# -----------------------------
# Source & YouTube enrichment (stubs with TODOs)
# -----------------------------
def find_journal_sources_for_week(queries: List[str]) -> List[Dict[str, Any]]:
    """
    TODO: Replace with Crossref/OpenAlex/Semantic Scholar.
    For now, LLM proposes 2-3 plausible sources. In production, verify & fetch pages from OA PDFs.
    """
    merged = "; ".join(queries[:3]) if queries else "general topic"
    j = llm_json(f"""
Propose 2-3 credible journal or conference sources with DOI/url for:
"{merged}"
Return: {{"sources":[{{"title":"","authors":[],"venue":"","year":2022,"url_or_doi":"","pages":"N/A"}}]}}
""")
    return j.get("sources", [])[:3]

def find_youtube_short_links(topics: List[str], max_videos: int = 2, max_minutes: int = 8) -> List[Dict[str, str]]:
    """
    TODO: Implement with YouTube Data API (ensure <= max_minutes).
    For now, ask LLM to suggest links. We trim to max_videos.
    """
    topic_str = "; ".join(topics[:3]) if topics else "course topic"
    j = llm_json(f"""
Suggest up to {max_videos} YouTube links (<= {max_minutes} minutes) relevant to: {topic_str}
Return: {{"videos":[{{"title":"","url":""}}]}}
""")
    vids = j.get("videos", [])[:max_videos]
    return vids

# -----------------------------
# Renderers: DOCX, PDF, PPTX, Questions
# -----------------------------
def add_logo_to_doc_header(document: docx.Document, logo_path: Optional[str]):
    if not logo_path:
        return
    section = document.sections[0]
    header = section.header
    p = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
    run = p.add_run()
    try:
        run.add_picture(logo_path, width=Inches(1.1))
        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    except Exception:
        pass

def generate_course_docx(plan: Dict[str, Any], out_path: str, logo_path: Optional[str]):
    doc = docx.Document()
    add_logo_to_doc_header(doc, logo_path)

    meta = plan.get("course_meta", {})
    title = meta.get("title") or "Course Guide"
    code = meta.get("code") or ""
    credits = meta.get("credits", 3)

    doc.add_heading(title, 0)
    if code:
        doc.add_paragraph(f"Course Code: {code}")
    doc.add_paragraph(f"Credits: {credits}")
    doc.add_paragraph(" ")

    for w in plan.get("weeks", []):
        doc.add_heading(f"Week {w.get('week')}: {w.get('title','')}", level=2)
        for label in ["topics", "objectives", "activities", "assessments"]:
            items = w.get(label, [])
            if items:
                doc.add_heading(label.capitalize(), level=3)
                for it in items:
                    doc.add_paragraph(it, style="List Bullet")
        sources = w.get("sources", [])
        if sources:
            doc.add_heading("Further Reading", level=3)
            for s in sources:
                line = f"{s.get('title','')} — {', '.join(s.get('authors',[]))} ({s.get('year','')}). {s.get('venue','')}. {s.get('url_or_doi','')} [pages: {s.get('pages','N/A')}]"
                doc.add_paragraph(line, style="List Number")
        vids = w.get("youtube", [])
        if vids:
            doc.add_heading("YouTube (≤8 min)", level=3)
            for v in vids:
                doc.add_paragraph(f"{v.get('title','Video')}: {v.get('url','')}", style="List Bullet")
        doc.add_paragraph(" ")

    doc.save(out_path)

def export_docx_to_pdf(docx_path: str, pdf_path: str):
    if DOCX2PDF_AVAILABLE:
        docx2pdf_convert(docx_path, pdf_path); return
    if REPORTLAB_AVAILABLE:
        d = docx.Document(docx_path)
        lines = [p.text for p in d.paragraphs]
        c = canvas.Canvas(pdf_path, pagesize=A4)
        width, height = A4
        x, y = 40, height - 50
        for line in lines:
            if y < 60:
                c.showPage(); y = height - 50
            c.drawString(x, y, line[:1200]); y -= 14
        c.save(); return
    raise RuntimeError("No PDF backend available. Install docx2pdf (Windows/macOS) or reportlab.")

def generate_pptx(plan: Dict[str, Any], out_path: str, logo_path: Optional[str], include_diagrams: bool = True):
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    meta = plan.get("course_meta", {})
    slide.shapes.title.text = meta.get("title") or "Course Deck"
    slide.placeholders[1].text = (meta.get("code") or "") + (" • Credits: " + str(meta.get("credits", 3)))
    if logo_path:
        try:
            slide.shapes.add_picture(logo_path, prs.slide_width - PPTInches(1.7), PPTInches(0.2), width=PPTInches(1.5))
        except Exception:
            pass

    # Overview
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Overview"
    tf = slide2.placeholders[1].text_frame
    tf.text = "This deck summarizes weekly topics, objectives, and resources."
    if logo_path:
        try:
            slide2.shapes.add_picture(logo_path, prs.slide_width - PPTInches(1.7), PPTInches(0.2), width=PPTInches(1.5))
        except Exception:
            pass

    # Per-week slides
    for w in plan.get("weeks", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Week {w.get('week')}: {w.get('title','')}"
        body = slide.placeholders[1].text_frame
        body.clear()
        body.paragraphs[0].text = "Topics:"
        for t in w.get("topics", [])[:5]:
            p = body.add_paragraph(); p.text = f"• {t}"; p.level = 1

        p = body.add_paragraph(); p.text = "Objectives:"; p.level = 0
        for t in w.get("objectives", [])[:5]:
            q = body.add_paragraph(); q.text = f"• {t}"; q.level = 1

        p = body.add_paragraph(); p.text = "Activities:"; p.level = 0
        for t in w.get("activities", [])[:5]:
            q = body.add_paragraph(); q.text = f"• {t}"; q.level = 1

        p = body.add_paragraph(); p.text = "Resources:"; p.level = 0
        for s in w.get("sources", [])[:3]:
            q = body.add_paragraph(); q.text = f"• {s.get('title','')} ({s.get('url_or_doi','')})"; q.level = 1
        for v in w.get("youtube", [])[:2]:
            q = body.add_paragraph(); q.text = f"• YouTube: {v.get('title','')} ({v.get('url','')})"; q.level = 1

        if include_diagrams:
            diagrams = set([d.lower() for d in w.get("diagram_suggestions", [])])
            if "pie" in diagrams: add_chart(slide, "pie")
            if "bar" in diagrams: add_chart(slide, "bar")
            if "flowchart" in diagrams: add_flowchart(slide)

        if logo_path:
            try:
                slide.shapes.add_picture(logo_path, prs.slide_width - PPTInches(1.7), PPTInches(0.2), width=PPTInches(1.5))
            except Exception:
                pass

    prs.save(out_path)

def add_chart(slide, chart_type="bar"):
    left = PPTInches(6.2); top = PPTInches(2.0); width = PPTInches(3.0); height = PPTInches(3.0)
    data = ChartData(); data.categories = ["A", "B", "C"]; data.add_series("Series 1", (2, 5, 3))
    if chart_type == "pie":
        slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, top, width, height, data)
    else:
        slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data)

def add_flowchart(slide):
    shapes = slide.shapes
    start = shapes.add_shape(MSO_SHAPE.OVAL, PPTInches(0.5), PPTInches(2.0), PPTInches(1.6), PPTInches(0.8)); start.text_frame.text = "Start"
    process = shapes.add_shape(MSO_SHAPE.RECTANGLE, PPTInches(2.4), PPTInches(2.0), PPTInches(2.2), PPTInches(0.9)); process.text_frame.text = "Process"
    decision = shapes.add_shape(MSO_SHAPE.DIAMOND, PPTInches(4.9), PPTInches(1.9), PPTInches(1.8), PPTInches(1.1)); decision.text_frame.text = "Decision?"
    end = shapes.add_shape(MSO_SHAPE.OVAL, PPTInches(7.0), PPTInches(2.0), PPTInches(1.6), PPTInches(0.8)); end.text_frame.text = "End"

def generate_questions_docx(plan: Dict[str, Any], out_path: str):
    d = docx.Document()
    d.add_heading("Question Bank", 0)
    j = llm_json(f"""
Create MCQs, short-answer, and essay questions aligned to the following course weeks.
Return JSON with keys mcq, short_answer, essay.
Weeks:
{json.dumps([{"week": w.get("week"), "title": w.get("title"), "topics": w.get("topics")} for w in plan.get("weeks",[])])}
""", temperature=0.5)

    def add_mcq(mcq):
        d.add_paragraph(f"Week {mcq.get('week')}: {mcq.get('question','')}")
        for idx, opt in enumerate(mcq.get("options", []), start=1):
            d.add_paragraph(f"   {idx}. {opt}")
        d.add_paragraph(f"Answer: {mcq.get('answer','')}")
        d.add_paragraph(" ")

    d.add_heading("Multiple Choice", level=1)
    for m in j.get("mcq", []): add_mcq(m)

    d.add_heading("Short Answer", level=1)
    for s in j.get("short_answer", []):
        d.add_paragraph(f"Week {s.get('week')}: {s.get('prompt','')}"); d.add_paragraph(" ")

    d.add_heading("Essay", level=1)
    for e in j.get("essay", []):
        d.add_paragraph(f"Week {e.get('week')}: {e.get('prompt','')}"); d.add_paragraph(" ")

    d.save(out_path)

# -----------------------------
# Orchestrator (respects flags)
# -----------------------------
def orchestrate(descriptor_text: str, instructions: str, weeks: int, logo_path: Optional[str], workdir: str, flags: Dict[str, Any]) -> Dict[str, str]:
    plan = plan_course_structure(descriptor_text, instructions, weeks)

    # Enrichment per flags
    for w in plan.get("weeks", []):
        if flags.get("include_sources", True):
            w["sources"] = find_journal_sources_for_week(w.get("source_queries", []))
        else:
            w["sources"] = []
        if flags.get("include_youtube", True):
            vids = find_youtube_short_links(
                w.get("topics", []),
                max_videos=int(flags.get("youtube_max", 2)),
                max_minutes=int(flags.get("youtube_max_minutes", 8)),
            )
            w["youtube"] = vids[: int(flags.get("youtube_max", 2))]
        else:
            w["youtube"] = []
        if not flags.get("include_diagrams", True):
            w["diagram_suggestions"] = []

    outputs: Dict[str, str] = {}

    # Always generate DOCX + JSON
    docx_path = os.path.join(workdir, "course_guide.docx")
    generate_course_docx(plan, docx_path, logo_path)
    outputs["course_guide_docx"] = docx_path

    raw_json_path = os.path.join(workdir, "plan.json")
    with open(raw_json_path, "w", encoding="utf-8") as f:
        json.dump(plan, f, indent=2, ensure_ascii=False)
    outputs["plan_json"] = raw_json_path

    # Optional PDF
    if flags.get("include_pdf", True):
        pdf_path = os.path.join(workdir, "course_guide.pdf")
        try:
            export_docx_to_pdf(docx_path, pdf_path)
            outputs["course_guide_pdf"] = pdf_path
        except Exception:
            # PDF backend not available; skip silently
            pass

    # Optional PPTX
    if flags.get("include_pptx", True):
        pptx_path = os.path.join(workdir, "course_deck.pptx")
        generate_pptx(plan, pptx_path, logo_path, include_diagrams=flags.get("include_diagrams", True))
        outputs["course_deck_pptx"] = pptx_path

    # Optional questions
    if flags.get("include_questions", True):
        qdocx_path = os.path.join(workdir, "questions.docx")
        generate_questions_docx(plan, qdocx_path)
        outputs["questions_docx"] = qdocx_path

    return outputs

# -----------------------------
# API
# -----------------------------
@router.post("/generate")
@router.post("/generate/")
async def generate(
    files: List[UploadFile] = File(..., description="Course descriptor(s) PDF/DOCX"),
    logo: Optional[UploadFile] = File(None, description="University logo PNG/JPG"),
    instructions: str = Form("", description="Extra instructions for tailoring"),
    weeks: int = Form(14),
    # Feature flags (strings in form-data; we'll parse to bools)
    include_pdf: str = Form("true"),
    include_pptx: str = Form("true"),
    include_questions: str = Form("true"),
    include_diagrams: str = Form("true"),
    include_sources: str = Form("true"),
    include_youtube: str = Form("true"),
    youtube_max: int = Form(2),
    youtube_max_minutes: int = Form(8),
):
    start = time.time()
    if weeks < 1 or weeks > 30:
        raise HTTPException(status_code=400, detail="weeks must be between 1 and 30")

    run_id = uuid.uuid4().hex[:10]
    workdir = os.path.abspath(os.path.join("runs", run_id))
    os.makedirs(workdir, exist_ok=True)

    try:
        # Extract text from uploads
        full_text = ""
        for f in files:
            name = f.filename or "upload"
            if not (name.lower().endswith(".pdf") or name.lower().endswith(".docx")):
                raise HTTPException(status_code=400, detail=f"Unsupported file: {name}")
            path = os.path.join(workdir, f"src_{name}")
            with open(path, "wb") as out:
                out.write(await f.read())
            full_text += extract_text_from_file(path) + "\n"

        if not full_text.strip():
            raise HTTPException(status_code=400, detail="No extractable text found in uploaded files.")

        logo_path = None
        if logo is not None and (logo.filename or "").lower().endswith((".png",".jpg",".jpeg")):
            lp = os.path.join(workdir, f"logo_{logo.filename}")
            with open(lp, "wb") as out:
                out.write(await logo.read())
            logo_path = lp

        # Parse flags
        flags = {
            "include_pdf": to_bool(include_pdf),
            "include_pptx": to_bool(include_pptx),
            "include_questions": to_bool(include_questions),
            "include_diagrams": to_bool(include_diagrams),
            "include_sources": to_bool(include_sources),
            "include_youtube": to_bool(include_youtube),
            "youtube_max": youtube_max,
            "youtube_max_minutes": youtube_max_minutes,
        }

        outputs = orchestrate(full_text, instructions, weeks, logo_path, workdir, flags)

        # Bundle ZIP
        zip_path = os.path.join(workdir, f"bundle_{run_id}.zip")
        with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as z:
            for k, p in outputs.items():
                if p and os.path.exists(p):
                    z.write(p, arcname=os.path.basename(p))

        elapsed = round(time.time() - start, 2)
        return {
            "status": "success",
            "duration_seconds": elapsed,
            "run_id": run_id,
            "download_zip": f"/download/{run_id}",
            "outputs": {k: os.path.basename(v) for k, v in outputs.items() if v}
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/download/{run_id}")
async def download_bundle(run_id: str):
    workdir = os.path.abspath(os.path.join("runs", run_id))
    if not os.path.isdir(workdir):
        raise HTTPException(status_code=404, detail="Run not found")
    zip_file = ""
    for fn in os.listdir(workdir):
        if fn.startswith("bundle_") and fn.endswith(".zip"):
            zip_file = os.path.join(workdir, fn); break
    if not zip_file or not os.path.exists(zip_file):
        raise HTTPException(status_code=404, detail="Bundle not found")
    return FileResponse(zip_file, filename=os.path.basename(zip_file), media_type="application/zip")

# -----------------------------
# App wiring
# -----------------------------
app.include_router(router)